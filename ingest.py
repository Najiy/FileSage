

import os
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from dotenv import load_dotenv
from docx import Document as DocxDocument

load_dotenv()

ROOT_DIR = '../DATA'
DB_DIR = 'db'
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
embeddings = OpenAIEmbeddings(model="text-embedding-3-small")

def extract_text_from_file(path):
    ext = path.lower().split('.')[-1]
    try:
        if ext in ['txt', 'md', 'js', 'py']:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == 'pdf':
            doc = fitz.open(path)
            return "\n".join(page.get_text() for page in doc)
        elif ext == 'pptx':
            prs = Presentation(path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == 'docx':
            doc = DocxDocument(path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif ext in ['png', 'jpg', 'jpeg']:
            return pytesseract.image_to_string(Image.open(path))
    except Exception as e:
        print(f"Error processing {path}: {e}")
    return ""

def ingest():
    docs = []
    for root, _, files in os.walk(ROOT_DIR):
        for file in files:
            path = os.path.join(root, file)
            ext = path.lower().split('.')[-1]
            print(f"Processing: {path}")
            text = extract_text_from_file(path)
            if not text.strip():
                continue

            if ext in ['png', 'jpg', 'jpeg']:
                # Image OCR result
                docs.append({
                    "text": text,
                    "metadata": {"source": path, "type": "image"}
                })
            else:
                # Normal doc
                chunks = text_splitter.split_text(text)
                for chunk in chunks:
                    docs.append({
                        "text": chunk,
                        "metadata": {"source": path, "type": "text"}
                    })

    db = FAISS.from_texts(
        texts=[doc["text"] for doc in docs],
        embedding=embeddings,
        metadatas=[doc["metadata"] for doc in docs]
    )
    db.save_local(DB_DIR)
    print("✅ Ingestion complete.")

if __name__ == "__main__":
    ingest()
